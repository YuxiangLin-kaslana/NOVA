#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/u/ylin30/sigLA")
FIG_DIR = ROOT / "slide_figures"
CASE_DIR = ROOT / "case_studies"
OUT = ROOT / "SigLA_experiment_slides.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

COLORS = {
    "dark": RGBColor(31, 45, 61),
    "gray": RGBColor(95, 108, 114),
    "blue": RGBColor(58, 110, 165),
    "teal": RGBColor(27, 153, 139),
    "green": RGBColor(87, 167, 115),
    "yellow": RGBColor(242, 193, 78),
    "red": RGBColor(217, 93, 57),
    "light": RGBColor(247, 249, 251),
    "white": RGBColor(255, 255, 255),
}


def set_bg(slide, color=COLORS["white"]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text: str, left, top, width, height, size=24, bold=False, color=COLORS["dark"], align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, title, Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.55), size=34, bold=True)
    if subtitle:
        add_text(slide, subtitle, Inches(0.58), Inches(1.02), Inches(12.0), Inches(0.35), size=16, color=COLORS["gray"])


def add_bullets(slide, bullets: list[str], left, top, width, height, size=21, color=COLORS["dark"]) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)


def add_card(slide, title: str, body: str, left, top, width, height, color) -> None:
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["light"]
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    add_text(slide, title, left + Inches(0.22), top + Inches(0.18), width - Inches(0.44), Inches(0.35), size=17, bold=True, color=color)
    add_text(slide, body, left + Inches(0.22), top + Inches(0.64), width - Inches(0.44), height - Inches(0.78), size=14, color=COLORS["dark"])


def add_full_image_slide(prs: Presentation, image_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide.shapes.add_picture(str(image_path), 0, 0, width=SLIDE_W, height=SLIDE_H)


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1: title.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, "SigLA Framework Validation", Inches(0.7), Inches(0.9), Inches(11.9), Inches(0.8), size=40, bold=True)
    add_text(
        slide,
        "Detector + policy training results and GPT agent behavior",
        Inches(0.74),
        Inches(1.78),
        Inches(11.5),
        Inches(0.4),
        size=20,
        color=COLORS["gray"],
    )
    add_card(slide, "Dataset", "SMD_1-7\n23,697 test points\n38 variables", Inches(0.78), Inches(3.0), Inches(2.7), Inches(1.55), COLORS["blue"])
    add_card(slide, "Training", "4,730 windows\n3,784 train\n946 validation", Inches(3.8), Inches(3.0), Inches(2.7), Inches(1.55), COLORS["teal"])
    add_card(slide, "Models", "MLPAnomalyDetector\nMLPActionPolicy", Inches(6.82), Inches(3.0), Inches(2.7), Inches(1.55), COLORS["green"])
    add_card(slide, "Agent", "100 GPT fallback calls\n100 valid decisions", Inches(9.84), Inches(3.0), Inches(2.7), Inches(1.55), COLORS["red"])
    add_text(slide, "/u/ylin30/sigLA", Inches(0.75), Inches(6.82), Inches(5), Inches(0.28), size=11, color=COLORS["gray"])

    # Slide 2: concise summary.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "What We Did", "Simple summary for slides")
    add_bullets(
        slide,
        [
            "Updated training to match the current detector -> concept -> policy -> agent framework.",
            "Trained detector and policy models with windows derived from the SMD_1-7 test split.",
            "Evaluated detector anomaly scoring, policy action prediction, and risk prediction.",
            "Ran a 100-window GPT agent test on the fallback pipeline.",
        ],
        Inches(0.85),
        Inches(1.7),
        Inches(5.9),
        Inches(4.7),
        size=22,
    )
    add_card(slide, "Detector", "Window ROC-AUC 0.840\nWindow AP 0.536\nOutputs continuous anomaly score", Inches(7.2), Inches(1.8), Inches(4.9), Inches(1.35), COLORS["blue"])
    add_card(slide, "Policy", "Action accuracy 0.989\nRisk F1 0.977\nMacro F1 on present actions 0.927", Inches(7.2), Inches(3.35), Inches(4.9), Inches(1.35), COLORS["green"])
    add_card(slide, "GPT Agent", "100 / 100 valid GPT decisions\nNo local fallback calls\nFollowed policy candidate actions", Inches(7.2), Inches(4.9), Inches(4.9), Inches(1.35), COLORS["red"])

    # Slides 3-7: generated figures.
    for image_name in [
        "01_framework_structure.png",
        "02_training_setup.png",
        "03_detector_results.png",
        "04_policy_results.png",
        "05_agent_results.png",
    ]:
        add_full_image_slide(prs, FIG_DIR / image_name)

    # Slides 8-9: case studies.
    for image_name in [
        "01_smd_detector_policy_case.png",
        "02_gpt_agent_fallback_case.png",
    ]:
        add_full_image_slide(prs, CASE_DIR / image_name)

    # Slide 10: takeaways.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Main Takeaways")
    add_card(
        slide,
        "Framework",
        "The code now trains components that match the current SigLA pipeline.",
        Inches(0.75),
        Inches(1.6),
        Inches(5.7),
        Inches(1.25),
        COLORS["blue"],
    )
    add_card(
        slide,
        "Policy",
        "The policy learns the weak action labels well in the test-derived experiment.",
        Inches(6.9),
        Inches(1.6),
        Inches(5.7),
        Inches(1.25),
        COLORS["green"],
    )
    add_card(
        slide,
        "Detector",
        "The detector is used as a score provider; downstream modules consume continuous reconstruction error.",
        Inches(0.75),
        Inches(3.15),
        Inches(5.7),
        Inches(1.25),
        COLORS["yellow"],
    )
    add_card(
        slide,
        "Agent",
        "GPT is operational and stable, but current context mostly leads it to explain the policy output.",
        Inches(6.9),
        Inches(3.15),
        Inches(5.7),
        Inches(1.25),
        COLORS["red"],
    )
    add_bullets(
        slide,
        [
            "Treat the current numbers as framework validation, not final benchmark results.",
            "Next: calibrate how policy/agent consume detector scores and run GPT agent on trained SMD pipeline outputs.",
        ],
        Inches(0.95),
        Inches(5.15),
        Inches(11.5),
        Inches(1.2),
        size=19,
        color=COLORS["gray"],
    )

    # Slide 11: paths.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Artifacts")
    add_bullets(
        slide,
        [
            "Detector run: /u/ylin30/sigLA/code/runs/detector_SMD_1-7_test_w50_s5",
            "Policy run: /u/ylin30/sigLA/code/runs/policy_SMD_1-7_test_w50_s5",
            "GPT fallback agent test: /u/ylin30/sigLA/code/runs/fallback_pipeline_gpt_test",
            "Case study figures: /u/ylin30/sigLA/case_studies",
            "Report: /u/ylin30/sigLA/SigLA_slide_report.md",
            "Figures: /u/ylin30/sigLA/slide_figures",
        ],
        Inches(0.85),
        Inches(1.7),
        Inches(11.8),
        Inches(4.5),
        size=18,
    )

    return prs


def main() -> None:
    prs = build_deck()
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
